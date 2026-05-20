from pptx import Presentation

def extract_text_from_ppt(ppt_path: str) -> str:
    prs = Presentation(ppt_path)
    all_text = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

        if slide_texts:
            all_text.append(f"[슬라이드 {slide_num}]")
            all_text.extend(slide_texts)

    return "\n".join(all_text)